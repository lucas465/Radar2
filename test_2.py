from numpy.ma.core import arctan2
from pandas.core.sorting import ensure_key_mapped
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
import pandas as pd
import datetime as dt
from pptx.enum.text import MSO_AUTO_SIZE
import os
import subprocess
import streamlit as st
###################################################
###################################################
#################BROWSER VERSION####################
###################################################
###################################################

# formula to calculate sizes: cm / 2.54 * 914400

##### streamlit interface setup
st.title("Excel File Uploader")

uploaded_file = st.file_uploader("Upload your Excel file", type=["xlsx", "xls"])

##################







#custom functions
def text_filter(input, target_list):
    final_herausgeber = []
    for i in input:
        r = i.split()
        tmp_list = []
        for j in r:
            test = [x for x in target_list if x in j]
            if test != []:
                tmp_list.append(test)
        if len(list(map(list, set(map(tuple,tmp_list))))) > 1:
           final_herausgeber.append("Double")
        elif tmp_list == []:
            final_herausgeber.append("tbd")
        else:
            final_herausgeber.append(tmp_list[0][0])
    return final_herausgeber
def text_filter_herausgeber(input, target_list):
    final_herausgeber = []
    for i in input:
        r = i.split()
        tmp_list = []
        for j in r:
            test = [x for x in target_list if x in j]
            if test != []:
                tmp_list.append(test)
        if tmp_list == []:
            final_herausgeber.append("tbd")
        else:
            final_herausgeber.append(tmp_list[0][0])
    return final_herausgeber
def single_word_filter(word,target_list):
        list_return = [1 for x in target_list if x in word]
        return list_return
def text_filter_light(input,target_list):
    final_herausgeber = []
    for i in input:
        r = i.split()
        tmp_list = []
        for j in r:
            test = [x for x in target_list if x in j]
            if test != []:
                tmp_list.append(test)
        if tmp_list == []:
            final_herausgeber.append("tbd")
        else:
            final_herausgeber.append(tmp_list)
    return final_herausgeber


# Base inputs
title_row = ["Regulierung","Art","Herausgeber","Inhalt","Veröffent-lichung","Relevant ab", "Status","Thema","Relevanz"] # each entry is 1 cell of the first row in the correct order
list_herausgeber = ["EBA","EZB","BIS","BCBS","ESRB","SRB","FSB","FMSG","ESRB","FATF","ESMA","EU-Amtsblatt","EU-Rat","EFAMA","EU-Kommission","EU-Rat","ESA","Accountancy Europe","EFRAG","EU-Parlament"]
list_art = ["DelVO","DVO","VO","Verordnung","verordnung","Gesetz","Pressemitteilung","Liste","Stellungnahme","Stellung","Tabelle","Mindeststandards","Statistik","Bericht","bericht","Leitlinien","Dashboard","Methodik","Richtlinie","Reporting Standards","Q&A","Call for Evidence","Urteil","RTS","ITS","Rundschreiben","Leitfaden","Memorandum of Understanding","Empfehlung","Vorläufige Einigung","Schreiben","Newsletter","Gesetz","Antwort","Positionspapier","Kommissionsvorschlag","Kommentar","Peer-Review","Klarstellung","Rede"]
list_regulation = ["CRR","CRD","BWG","Taxonomie-VO","Taxonomie Verordnung","MIFID","MiFID","MiFIR","CSDR","EMIR","CCP","DORA","MiCAR","MiCA","BRRD","SRM","SRM-VO","SEPA-VO","SEPA","PSD2","PSD","CSDDD","CSRD","EuGBS-VO","EuGBS","Verbriefungs-VO","UCITS","AIFMD","Prospekt-VO","MMF-VO"]
list_entwurf = ["Entwurf","Konsultation"]
map_granular_category_1 = {"CRR": "CRR/CRD/BWG","CRD": "CRR/CRD/BWG", "BWG": "CRR/CRD/BWG","CSDR": "Asset Management", "AIFMD": "Asset Management",
                             "UCITS": "Asset Management", "OGAW": "Asset Management","InvestIFG": "Asset Management","CCP": "Asset Management","EMIR": "Asset Management", "MiFID": "MiFID/MiFIR/WAG",
                             "MiFIR": "MiFID/MiFIR/WAG", "WAG": "MiFID/MiFIR/WAG", "PRIIPS": "PRIIPS","IFR": "IFR/IFD/WPFG", "IFD": "IFR/IFD/WPFG", "WPFG": "IFR/IFD/WPFG",
                             "DORA": "DORA","MICAR": "MICAR", "MICA":"MICAR","MiCAR":"MICAR","MiCA":"MICAR"}
################################################

#searching for input file in folder
path = os.path.abspath(os.getcwd()) # getting path of working directory
dir_list = os.listdir(path) # every file in the path

 #################################

### extracting content form Input table
df_news = pd.read_excel(uploaded_file)   # reading excel with pandas
inhalt = [x.replace("\n","") for x in df_news.iloc[:,1].tolist()] # reading inhalt from second column and filtering line breaks
teaser = [x.replace("\n","") for x in df_news.iloc[:,5].tolist()] # reading inhalt from second column and filtering line breaks
status = ["test"]*len(inhalt) # temporary filler
relevanz = ["niedrig"]*len(inhalt) # temporary filler
thema = df_news.iloc[:, 0].tolist() # reading thema/title from column 1 pf the input file
###########################################


#### checking every "inhalt" input (inhalt, thema, teaser) for regulierung
regulierung1 = text_filter(thema,list_regulation)
regulierung2 = text_filter(inhalt,list_regulation)
regulierung3 = text_filter(teaser,list_regulation)
regulierung = []
for d in range(0,len(regulierung1)):
    if regulierung1[d] == "tbd" and regulierung2[d] == "tbd" and regulierung3[d] == "tbd":
        regulierung.append("keine")
    elif regulierung1[d] != "tbd":
        regulierung.append(regulierung1[d])
        pass
    elif regulierung2[d] != "tbd":
        regulierung.append(regulierung2[d])
        pass
    elif regulierung3[d] != "tbd":
        regulierung.append(regulierung3[d])
        pass


####### list prep for the granular mapping of categories for later (getting doubles out of the way, if they´re of the same category)
gran_reg = []
regulierung_filter_light = text_filter_light(thema,list_regulation)
flat_reg_filt_light = [item for sublist in regulierung_filter_light for item in (sublist if isinstance(sublist, list) else [sublist])]
print(flat_reg_filt_light)
print(regulierung_filter_light)
for f in flat_reg_filt_light:
    if f == "tbd":
        gran_reg.append("tbd")
        pass
    elif len(f) == 1:
        gran_reg.append(map_granular_category_1.get(f[0]))
    else:
        tmp = []
        for h in f:
            tmp.append(map_granular_category_1.get(h))
        if len(set(tmp)) == 1:
            gran_reg.append((list((set(tmp)))[0]))
        else:
            gran_reg.append("multiple regulations")
print(len(flat_reg_filt_light))
print(len(gran_reg))
print(len(regulierung_filter_light))
print(regulierung_filter_light)


################################################

#######################################################################
#### pulishing date
veröffentlichung = [x.strftime("%Y-%m-%d") for x in df_news.iloc[:,2].tolist()] #reading publishing date from column 3 of the input table
relevant_ab = veröffentlichung # start + end date the same for now
###########################################

#### checking every "inhalt" input (inhalt, thema, teaser) for art
art1 = text_filter_herausgeber(thema, list_art)
art2 = text_filter_herausgeber(teaser, list_art)
art = []
for d in range(0,len(art1)):
    if art1[d] == "tbd" and art2[d] == "tbd":
        art.append("tbd")
    elif art1[d] != "tbd":
        art.append(art1[d])
        pass
    elif art2[d] != "tbd":
        art.append(art2[d])
        pass
art = ["Bericht" if x == "bericht" else x for x in art]
###################################################################################
############### konsolidierung art
art = ["Bericht" if x == "Positionspapier" else x for x in art]
art = ["Bericht" if x == "Newsletter" else x for x in art]
art = ["Bericht" if x == "Newsletter" else x for x in art]

###################################




###################################



###################################





#### filtering for herausgeber
herausgeber = text_filter_herausgeber(thema,list_herausgeber)


### category sorting based on given categories:
de_categories_map = {"Digitalisierung & IT-Technologie":"Digitalisierung", "Eigenmittel & Liquidität": "Bankenrecht","ESG": "Sustainable Finance",
                     "Finanzmärkte": "Wertpapieraufsichtsrecht", "Geldwäschebekämpfung": "AML und Sanktionen","Künstliche Intelligenz": "Digitalisierung", "Meldewesen & Offenlegung": "tbd","n/a": "tbd","Risikomanagement & Governance": "tbd","Sanierung & Abwicklung": "Bankenrecht","Sonstiges":"tbd",
                     "Verbraucherschutz": "tbd","Zahlungsverkehr":"Bankenrecht" }

df_news["Category1"] = df_news["Kategorien"].map(de_categories_map)
Category1 = df_news["Category1"].tolist()
###############################################################
#mapping based on sector:
de_industry_map = {"Asset Management":"Wertpapieraufsichtsrecht","Banking":"Bankenrecht","FS Allgemein": "tbd","Insurance":"tbd"}
df_news["Category3"] = df_news["Industry"].map(de_industry_map)
Category3 = df_news["Category3"].tolist()
##############################################

#######################################################################################
### category variable for later sorting
Category2 = []
for w in regulierung:
    if w == "tbd":
        Category2.append("tbd")
    elif w == "CRR" or w == "CRD" or w == "BWG" or w == "SRM" or w == "SEPA" or w == "PSD" or w == "PSD2" or w == "BRRD":
        Category2.append("Bankenrecht")
    elif w == "MIFID" or w == "MiFID" or w == "MiFIR" or w == "MIFIR" or w == "CSDR" or w == "EMIR" or w == "CCP":
        Category2.append("Wertpapieraufsichtsrecht")
    elif w == "DORA" or w == "MiCA" or w == "MiCAR":
        Category2.append("Digitalisierung")
    elif w == "Taxonomie-VO" or w == "CSDDD" or w == "CSRD" or w == "EuGBS":
        Category2.append("Sustainable Finance")
    else:
        Category2.append("???")
#########################
## Entwurf final
list_entwurf = text_filter(inhalt, list_entwurf)
print(list_entwurf)
entwurf_final = []
for j in list_entwurf:
    if j == "Entwurf":
        entwurf_final.append("Entwurf")
    elif j == "Konsultation":
        entwurf_final.append("Konsultation")
    else:
        entwurf_final.append("Final")
###################
############ if "art" comes before "herausgeber" the article get´s filtered (since there must be a non-relevant "herausgeber"))
indexes_1 = []
indexes_2 = []
for t in range(0,len(thema)):
    tmp_indexes_art = []
    tmp_indexes_herausgeber = []
    q = thema[t].split()
    for g in q:
        tmp_indexes_art.append(single_word_filter(g,list_art))
        tmp_indexes_herausgeber.append(single_word_filter(g,list_herausgeber))
        sum_flat_list = sum([x for sublist in tmp_indexes_herausgeber for x in sublist])
    for j in range(0,len(tmp_indexes_art)):
        if sum_flat_list == 0:
            break
        elif tmp_indexes_art[j] != [] and tmp_indexes_herausgeber[j] == []:
            indexes_1.append(t)
            break
        elif tmp_indexes_art[j] != [] and tmp_indexes_herausgeber[j] != []:
            break
        elif tmp_indexes_art[j] == [] and tmp_indexes_herausgeber[j] != []:
            break
    tmp_indexes_art = []
    tmp_indexes_herausgeber = []
    m = teaser[t].split()
    for g in m:
        tmp_indexes_art.append(single_word_filter(g, list_art))
        tmp_indexes_herausgeber.append(single_word_filter(g, list_herausgeber))
        sum_flat_list = sum([x for sublist in tmp_indexes_herausgeber for x in sublist])
    for j in range(0, len(tmp_indexes_art)):
        if sum_flat_list == 0:
            break
        elif tmp_indexes_art[j] != [] and tmp_indexes_herausgeber[j] == []:
            indexes_2.append(t)
            break
        elif tmp_indexes_art[j] != [] and tmp_indexes_herausgeber[j] != []:
            break
        elif tmp_indexes_art[j] == [] and tmp_indexes_herausgeber[j] != []:
            break

################################################################################

#### filter non-relevant herausgeber
indexes_3 = []
for u in range(0,len(herausgeber)):
    if herausgeber[u] == "tbd" or art1[u] == "Double":
        indexes_3.append(u)
indexes_4 = []
for j in range(0,len(art)):
    if j == "Rede":
        indexes_4.append(j)


indexes = list(set(indexes_1 + indexes_2 + indexes_3 + indexes_4))

filtered_regulierung = [item for i, item in enumerate(regulierung) if i not in indexes]
filtered_art = [item for i, item in enumerate(art) if i not in indexes]
filtered_herausgeber = [item for i, item in enumerate(herausgeber) if i not in indexes]
filtered_inhalt = [item for i, item in enumerate(inhalt) if i not in indexes]
filtered_veröffentlichung = [item for i, item in enumerate(veröffentlichung) if i not in indexes]
filtered_relevant_ab = [item for i, item in enumerate(relevant_ab) if i not in indexes]
filtered_status = [item for i, item in enumerate(status) if i not in indexes]
filtered_thema = [item for i, item in enumerate(thema) if i not in indexes]
filtered_relevanz = [item for i, item in enumerate(relevanz) if i not in indexes]
filtered_category1 = [item for i, item in enumerate(Category1) if i not in indexes]
filtered_category2 = [item for i, item in enumerate(Category2) if i not in indexes]
filtered_entwurf = [item for i, item in enumerate(entwurf_final) if i not in indexes]
filtered_category3 = [item for i, item in enumerate(Category3) if i not in indexes]
filtered_teaser = [item for i, item in enumerate(teaser) if i not in indexes]
filtered_gran_reg = [item for i, item in enumerate(gran_reg) if i not in indexes]

##########################################################################################




##########################################################################################






############ determining final high level category (Bankenrecht, Werpapierrecht etc.)
category_final = []
for k in range(0,len(filtered_category1)):
    if filtered_category3[k] != "tbd" and filtered_category3[k] != "???":
        category_final.append(filtered_category3[k])
        pass
    elif filtered_category2[k] != "tbd" and filtered_category2[k] != "???":
        category_final.append(filtered_category2[k])
        pass
    elif filtered_category1[k] != "tbd" and filtered_category1[k] != "???":
        category_final.append(filtered_category1[k])
        pass
    else:
        category_final.append("tbd")
##############################################
#### Determining final granular categories
gran_categ_filt = []

for k in range(0,len(category_final)):
    if category_final[k] == "Bankenrecht":
        g = "Weitere Neuheiten Bankenrecht"
    elif category_final[k] == " Wertpapieraufsichtsrecht":
        g = "Weitere Neuheiten Wertpapieraufsichtsrecht"
    elif category_final[k] == "Digitalisierung":
        g = "Weitere Neuheiten Digitalisierung"
    map_granular_category = {"CRR": "CRR/CRD/BWG","CRD": "CRR/CRD/BWG", "BWG": "CRR/CRD/BWG","keine":g ,"CSDR": "Asset Management", "AIFMD": "Asset Management",
                             "UCITS": "Asset Management", "OGAW": "Asset Management","InvestIFG": "Asset Management","CCP": "Asset Management","EMIR": "Asset Management", "MiFID": "MiFID/MiFIR/WAG",
                             "MiFIR": "MiFID/MiFIR/WAG", "WAG": "MiFID/MiFIR/WAG", "PRIIPS": "PRIIPS","IFR": "IFR/IFD/WPFG", "IFD": "IFR/IFD/WPFG", "WPFG": "IFR/IFD/WPFG",
                             "DORA": "DORA","MICAR": "MICAR", "MICA":"MICAR","MiCAR":"MICAR","MiCA":"MICAR"}

    gran_categ_filt.append(map_granular_category.get(filtered_regulierung[k]))

fin_gran_categ = []
for j in range(0,len(gran_categ_filt)):
    if gran_categ_filt[j] == gran_reg[j]:
        fin_gran_categ.append(gran_categ_filt[j])
    elif gran_reg[j] == "tbd"  or gran_reg[j] == None and gran_categ_filt[j] != None:
        fin_gran_categ.append(gran_categ_filt[j])
    elif gran_reg[j] == None and gran_reg[j] != None and gran_categ_filt[j] != "tbd":
        fin_gran_categ.append(gran_categ_filt[j])
print(len(gran_categ_filt))
print(len(filtered_gran_reg))
print(len(filtered_relevant_ab))
#####################

################## create dataframe with final inputs for the ppt

dict_fin_input = {"Regulierung":filtered_regulierung, "Art":filtered_art, "Herausgeber":filtered_herausgeber,"Inhalt":filtered_teaser,"Veröffentlichung":filtered_veröffentlichung,
                  "Veröffentlichung":filtered_veröffentlichung,"Relevant":filtered_relevant_ab,"Status":filtered_status,"Thema":filtered_thema,"Relevanz":filtered_relevanz,
                  "Category": category_final,"Status":filtered_entwurf}
df_input_clean = pd.DataFrame(dict_fin_input)



### highlighting rows in excel that need human attention
def highlight_cell(s):
    is_equal = s.isin(["Double","tbd",""]) | s.isna()
    return ['background-color: yellow' if v else '' for v in is_equal]

df_input_clean.style.apply(highlight_cell).to_excel("final_inputs.xlsx")
#####

if os.name == 'nt':  # Windows
    subprocess.Popen(['start', "final_inputs.xlsx"], shell=True)
input("Press Enter after editing and closing the Excel file...")

#############################################################

# sorting final dataframe and adding category numbers for sorting
df_appended = pd.read_excel("final_inputs.xlsx")
category_map = {"Bankenrecht": int(1), "Sustainable Finance": int(2),"AML und Sanktionen":int(3), "Wertpapieraufsichtsrecht": int(4), "Digitalisierung": int(5), "tbd":int(6),"???":int(6)}
df_appended["category_number"] = df_appended["Category"].map(category_map)
sorted_df = df_appended.sort_values(by=["category_number","Veröffentlichung"])
sorted_df = sorted_df.iloc[: , 1:]
chapters = sorted_df["category_number"].tolist()
# number of slides and content distribution
adder = 0 #counts up form 0 in steps of 1 and accounts for the shift that using 2 articles on page causes
slides = [] # list to store the number of articles on each slide

inhalt_sort = sorted_df["Inhalt"].tolist()
for k in range(0,len(inhalt_sort)):
    k = k + adder
    if k < len(inhalt_sort) -1  and len(inhalt_sort[k]) + len(inhalt_sort[k+1]) > 1050 or k < len(inhalt_sort) -1 and chapters[k] != chapters[k+1]:
        slides.append(1)
    elif k < len(inhalt_sort)-1:
        slides.append(2)
        adder = adder + 1
    elif sum(slides) == len(inhalt_sort):
        break
    else:
        slides.append(1)
        break

#############################################################
article_counter = 0 # counts the number of articles already included
print(sorted_df)
# Initializing Presentation
prs = Presentation()
###########################
for num in range(0,len(slides)):
    print(article_counter)
    print(sorted_df["Category"].iloc[article_counter])
    #Setting Slide size to the standard KPMG Widescreen Size
    prs.slide_width = 12240000
    prs.slide_height = 7200000
    ########################################################

    # Defining Slide layout (0 just title, 1 title and content etc.) and creating the slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    #######################################################################################
    #adding slide title to each slide
    left = 908559
    top = 540000
    width = 8590000
    height = 6000000
    text_box_title = slide.shapes.add_textbox(left,top,width,height)
    title_frame = text_box_title.text_frame
    title_par = title_frame.paragraphs[0]
    title_par.alingment = PP_ALIGN.LEFT
    run_title = title_par.add_run()
    if sorted_df["Category"].iloc[article_counter] == "Bankenrecht":
        run_title.text = "Aktuelle Entwicklungen iZm Bankenrecht"
    elif sorted_df["Category"].iloc[article_counter] == "Sustainable Finance":
        run_title.text = "Aktuelle Entwicklungen iZm Sustainable Finance"
    elif sorted_df["Category"].iloc[article_counter] == "Wertpapieraufsichtsrecht":
        run_title.text = "Aktuelle Entwicklungen iZm Wertpapieraufsichtsrecht"
    elif sorted_df["Category"].iloc[article_counter] == "Digitalisierung":
        run_title.text = "Aktuelle Entwicklungen iZm Digitalisierung"
    elif sorted_df["Category"].iloc[article_counter] == "AML und Sanktionen":
        run_title.text = "Aktuelle Entwicklungen iZm AML und Sanktionen"
    else:
        print("Category missing")
        exit()

    run_title.font.name = "KPMG Bold"
    run_title.font.size = Pt(36)
    run_title.font.color.rgb = RGBColor(0,51,151)

    ###############
    ######## Supertitle for each Slide ###########
    left_st = 908559
    top_st = 370000
    width_st = 8590000
    height_st = 6000000
    text_box_st = slide.shapes.add_textbox(left_st,top_st,width_st,height_st)
    st_frame = text_box_st.text_frame
    st_par = st_frame.paragraphs[0]
    st_par.alingment = PP_ALIGN.LEFT
    run_st = st_par.add_run()
    run_st.text = sorted_df["Category"].iloc[article_counter]
    run_st.font.name = "Arial"
    run_st.font.size = Pt(14)
    run_st.font.color.rgb = RGBColor(0, 51, 151)
    run_st.font.bold = True

    ####################################



    ########################################

    #adding table ####################################
    left = 1007559   # space from left edge of slide
    top = 1440000    # space from the top if the slide
    width = 7980000 # width of table evenly distirbuted
    height = 6000000 # height evenly distributed
    shape = slide.shapes.add_table(slides[num]+1,9,left,top,width,height)
    table = shape.table
    table.rows[0].height = 450000
    table.columns[3].width = 2610000
    table.columns[7].width = 1440000
    if slides[num] == 1:
        table.rows[1].height = 100000
    else:
        table.rows[1].height = 100000
        table.rows[2].height = 100000
    ################################################

    # creating top row with text and fill (standard for all slides)
    for i in range(0,9):
        temp_cell = table.cell(0,i) # looping through each cell of the first row
        temp_cell.vertical_anchor = MSO_ANCHOR.MIDDLE   #aligning text in the middle of the cell verticaly

        fill_temp = temp_cell.fill  #colouring table cell
        fill_temp.solid()
        if sorted_df["Category"].iloc[article_counter] == "Bankenrecht":
            fill_temp.fore_color.rgb = RGBColor(39,51,141) #kpmg blue
        elif sorted_df["Category"].iloc[article_counter] == "Sustainable Finance":
            fill_temp.fore_color.rgb = RGBColor(30,73,226) #cobalt blue
        elif sorted_df["Category"].iloc[article_counter] == "Wertpapieraufsichtsrecht":
            fill_temp.fore_color.rgb = RGBColor(171,13,130) #pink
        elif sorted_df["Category"].iloc[article_counter] == "Digitalisierung":
            fill_temp.fore_color.rgb = RGBColor(9,142,126) #green
        elif sorted_df["Category"].iloc[article_counter] == "AML und Sanktionen":
            fill_temp.fore_color.rgb = RGBColor(81,13,187) #purple

        text_frame_tmp = temp_cell.text_frame # creating text frame (will contain text later)
        text_frame_tmp.word_wrap = True # turning word wrap on for each cell

        temp_par = text_frame_tmp.paragraphs[0] #creating paragraph inside the text frame
        temp_par.alignment = PP_ALIGN.CENTER # aligning the paragraph --> horizontal alignment

        run = temp_par.add_run()  # creating run (contains the actual text)
        run.text = title_row[i] # looping through list with standard text for first row (same on all slides)
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.bold = True
    article_counter = article_counter + slides[num]
    ###################################################

    # filling rows (fill colour, content)
    for k in range(0,9):
        for l in range(1,slides[num]+1):
            cell_content = table.cell(l,k)
            cell_content.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill_content = cell_content.fill
            fill_content.solid()
            if l == 1:
                fill_content.fore_color.rgb = RGBColor(229,229,229) # darker grey
            else:
                fill_content.fore_color.rgb = RGBColor(242,242,242) # lighter grey

            text_frame_content = cell_content.text_frame
            text_frame_content.word_wrap = True

            par_content = text_frame_content.paragraphs[0]
            par_content.alignment = PP_ALIGN.CENTER

            run_cont = par_content.add_run()
            run_cont.text = sorted_df.iloc[l-1+sum(slides[0:num])].iloc[k]
            run_cont.font.name = "Arial"
            run_cont.font.size = Pt(9)
            run_cont.font.color.rgb = RGBColor(0,51,141)
    ####################################################








# Saving presentation (creates a new presentation or overwrites the
# presentation in the folder where the script is saved)
prs.save("test2.pptx")
####################################################################