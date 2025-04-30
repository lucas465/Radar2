from pandas.core.sorting import ensure_key_mapped
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
import pandas as pd
import datetime as dt
from pptx.enum.text import MSO_AUTO_SIZE
import os

#custom functions
def text_filter(input, target_list):
    final_herausgeber = []
    for i in input:
        r = i.split()
        tmp_list = []
        for j in r:
            test = [x for x in target_list if x in j]
            if test != []:
                tmp_list.append(test)
        if len(tmp_list) > 1:
           final_herausgeber.append("Double")
        elif tmp_list == []:
            final_herausgeber.append("tbd")
        else:
            final_herausgeber.append(tmp_list[0][0])
    return final_herausgeber

def double_filter(herausgeber,art,word):
    herausgeber_fin = []
    for i in range(0, len(herausgeber)):
        if art[i] == word or herausgeber[i] == word:
            herausgeber_fin.append("manual entry")
        else:
            herausgeber_fin.append(herausgeber[i])
    return  herausgeber_fin



# formula to calculate sizes: cm / 2.54 * 914400

# Base inputs
title_row = ["Regulierung","Art","Herausgeber","Inhalt","Veröffent-lichung","Relevant ab", "Status","Thema","Relevanz"] # each entry is 1 cell of the first row in the correct order
list_herausgeber = ["EBA","ECB","EZB","BIS","BCBS","SRB","FSB","FMSG","ESRB","FATF","ESMA","EU-Amtsblatt","EU-Rat","EFAMA","EIOPA","EU-Kommission","EU-Rat","ESA","Accountancy Europe","EFRAG","EU-Parlament"]
list_art = ["DelVO","VO","DVO","Gesetz","Pressemitteilung","Liste","Stellungnahme","Tabelle","Mindeststandards","Statistik","Bericht","bericht","Leitlinien","Dashboard","Methodik","Richtlinie","Reporting Standards","Q&A","Call for Evidence","Urteil","RTS","ITS","Rundschreiben","Leitfaden","Memorandum of Understanding","Empfehlung","Vorläufige Einigung","Schreiben","Newsletter","Gesetz","Antwort","Stellung","Vorschlag","Vorschläge","Positionspapier","Kommissionsvorschlag","Konsultation","Kommentar","Peer-Review","Klarstellung","Daten","Rede"]
################################################
#searching for input file in folder
path = os.path.abspath(os.getcwd())
print(path)
#################################

### extracting content form DE Table
df_de_news = pd.read_excel("data_2.xlsx")
inhalt = [x.replace("\n","") for x in df_de_news["Textkörper"].tolist()]
regulierung = ["test"]*len(inhalt)   # temporary filler
status = ["test"]*len(inhalt) # temporary filler
relevanz = ["test"]*len(inhalt) # temporary filler
thema = df_de_news["Title"].tolist()
veröffentlichung = [x.strftime("%Y-%m-%d") for x in df_de_news["Issue Date"].tolist()]
relevant_ab = veröffentlichung
art = text_filter(thema, list_art)
herausgeber = double_filter(text_filter(thema,list_herausgeber),art,"Double")

#### filter non-relevant herausgeber
indexes = []
for u in range(0,len(herausgeber)):
    if herausgeber[u] == "tbd" or herausgeber[u] == "manual entry":
        indexes.append(u)
filtered_regulierung = [item for i, item in enumerate(regulierung) if i not in indexes]
filtered_art = [item for i, item in enumerate(art) if i not in indexes]
filtered_herausgeber = [item for i, item in enumerate(herausgeber) if i not in indexes]
filtered_inhalt = [item for i, item in enumerate(inhalt) if i not in indexes]
filtered_veröffentlichung = [item for i, item in enumerate(veröffentlichung) if i not in indexes]
filtered_relevant_ab = [item for i, item in enumerate(relevant_ab) if i not in indexes]
filtered_status = [item for i, item in enumerate(status) if i not in indexes]
filtered_thema = [item for i, item in enumerate(thema) if i not in indexes]
filtered_relevanz = [item for i, item in enumerate(relevanz) if i not in indexes]

################## create dataframe with final inputs for the ppt
dict_fin_input = {"Regulierung":filtered_regulierung, "Art":filtered_art, "Herausgeber":filtered_herausgeber,"Inhalt":filtered_inhalt,"Veröffentlichung":filtered_veröffentlichung,"Veröffentlichung":filtered_veröffentlichung,"Relevant":filtered_relevant_ab,"Status":filtered_status,"Thema":filtered_thema,"Relevanz":filtered_relevanz}
df_input_clean = pd.DataFrame(dict_fin_input)
print(df_input_clean)
#############################################################

# number of slides and content distribution
adder = 0 #counts up form 0 in steps of 1 and accounts for the shift that using 2 articles on page causes
slides = [] # list to store the number of articles on each slide
for k in range(0,len(filtered_inhalt)):
    k = k + adder
    if k < len(filtered_inhalt) -1  and len(filtered_inhalt[k]) + len(filtered_inhalt[k+1]) > 1050:
        slides.append(1)
    elif k < len(filtered_inhalt)-1:
        slides.append(2)
        adder = adder + 1
    elif sum(slides) == len(filtered_inhalt):
        break
    else:
        slides.append(1)
        break

#############################################################

# Initializing Presentation
prs = Presentation()
###########################
for num in range(0,len(slides)):
    #Setting Slide size to the standard KPMG Widescreen Size
    prs.slide_width = 12240000
    prs.slide_height = 7200000
    ########################################################

    # Defining Slide layout (0 just title, 1 title and content etc.) and creating the slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    #######################################################################################
    #adding slide title to each slide
    left = 908559
    top = 540000
    width = 8590000
    height = 6000000
    text_box_title = slide.shapes.add_textbox(left,top,width,height)
    title_frame = text_box_title.text_frame
    title_par = title_frame.paragraphs[0]
    title_par.alingment = PP_ALIGN.LEFT
    run_title = title_par.add_run()
    run_title.text = "Regulatory Radar"
    run_title.font.name = "KPMG Bold"
    run_title.font.size = Pt(36)
    run_title.font.color.rgb = RGBColor(0,51,151)
    ########################################

    #adding table ####################################
    left = 1007559   # space from left edge of slide
    top = 1440000    # space from the top if the slide
    width = 8590000 # width of table evenly distirbuted
    height = 6000000 # height evenly distributed
    shape = slide.shapes.add_table(slides[num]+1,9,left,top,width,height)
    table = shape.table
    table.rows[0].height = 450000
    table.columns[3].width = 2610000
    if slides[num] == 1:
        table.rows[1].height = 100000
    else:
        table.rows[1].height = 100000
        table.rows[2].height = 100000
    ################################################

    # creating top row with text and fill (standard for all slides)
    for i in range(0,9):
        temp_cell = table.cell(0,i) # looping through each cell of the first row
        temp_cell.vertical_anchor = MSO_ANCHOR.MIDDLE   #aligning text in the middle of the cell verticaly

        fill_temp = temp_cell.fill  #colouring table cell
        fill_temp.solid()
        fill_temp.fore_color.rgb = RGBColor(0,51,141) #kpmg blue

        text_frame_tmp = temp_cell.text_frame # creating text frame (will contain text later)
        text_frame_tmp.word_wrap = True # turning word wrap on for each cell

        temp_par = text_frame_tmp.paragraphs[0] #creating paragraph inside the text frame
        temp_par.alignment = PP_ALIGN.CENTER # aligning the paragraph --> horizontal alignment

        run = temp_par.add_run()  # creating run (contains the actual text)
        run.text = title_row[i] # looping through list with standard text for first row (same on all slides)
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.bold = True
    ###################################################

    # filling rows (fill colour, content)
    for k in range(0,9):
        for l in range(1,slides[num]+1):
            cell_content = table.cell(l,k)
            cell_content.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill_content = cell_content.fill
            fill_content.solid()
            if l == 1:
                fill_content.fore_color.rgb = RGBColor(229,229,229) # darker grey
            else:
                fill_content.fore_color.rgb = RGBColor(242,242,242) # lighter grey

            text_frame_content = cell_content.text_frame
            text_frame_content.word_wrap = True

            par_content = text_frame_content.paragraphs[0]
            par_content.alignment = PP_ALIGN.CENTER

            run_cont = par_content.add_run()
            run_cont.text = df_input_clean.iloc[l-1+sum(slides[0:num])].iloc[k]
            run_cont.font.name = "Arial"
            run_cont.font.size = Pt(9)
            run_cont.font.color.rgb = RGBColor(0,51,141)




    ####################################################








# Saving presentation (creates a new presentation or overwrites the
# presentation in the folder where the script is saved)
prs.save("test2.pptx")
####################################################################